"""
==========================================================
  ZENITH OX â€” Secure Intelligent Research Assistant
  Backend: Flask + Groq (fast cloud LLMs) + TF-IDF Vector Memory
  Auth:    Email+Password  OR  Google OAuth  (or both)
  Web:     Tavily API via raw `requests` (no SDK)
  Stage 2: AI Mode Selection (6 modes)
  
  UPDATED: Gemini -> Groq cloud API via OpenAI-compatible endpoint
==========================================================
"""

import os
import json
import secrets
import re
import time
import zipfile
import io
from datetime import datetime

import requests
import numpy as np

from openai import OpenAI
from flask import (
    Flask, render_template, request, jsonify,
    redirect, url_for, session, flash, send_from_directory,
)
from werkzeug.security import generate_password_hash, check_password_hash
from authlib.integrations.flask_client import OAuth
from dotenv import load_dotenv
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

# File analysis imports (universal file support)
import csv
import base64
import mimetypes

# ----------------------------------------------------------
# 1. ENVIRONMENT
# ----------------------------------------------------------
load_dotenv()

BASE_DIR    = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

USERS_FILE  = os.path.join(BASE_DIR, "users.json")
MEMORY_FILE = os.path.join(BASE_DIR, "memory.json")

WRITABLE_USERS  = "/tmp/users.json"
WRITABLE_MEMORY = "/tmp/memory.json"

GROQ_API_KEY        = os.getenv("GROQ_API_KEY", "")
TAVILY_API_KEY      = os.getenv("TAVILY_API_KEY", "")
SECRET_KEY          = os.getenv("FLASK_SECRET_KEY", secrets.token_hex(32))

GOOGLE_CLIENT_ID     = os.getenv("GOOGLE_CLIENT_ID", "")
GOOGLE_CLIENT_SECRET = os.getenv("GOOGLE_CLIENT_SECRET", "")

groq_client = OpenAI(
    api_key=GROQ_API_KEY,
    base_url="https://api.groq.com/openai/v1",
)

# ----------------------------------------------------------
# 2. FLASK APP
# ----------------------------------------------------------
app = Flask(
    __name__,
    template_folder=os.path.join(BASE_DIR, "templates"),
    static_folder=os.path.join(BASE_DIR, "static"),
)
app.secret_key = SECRET_KEY

MEMORY_LIMIT = 15
TOP_K_MEMORY = 3

# ----------------------------------------------------------
# 2b. GOOGLE OAUTH SETUP
# ----------------------------------------------------------
oauth = OAuth(app)
google = None
if GOOGLE_CLIENT_ID and GOOGLE_CLIENT_SECRET:
    google = oauth.register(
        name="google",
        client_id=GOOGLE_CLIENT_ID,
        client_secret=GOOGLE_CLIENT_SECRET,
        server_metadata_url="https://accounts.google.com/.well-known/openid-configuration",
        client_kwargs={"scope": "openid email profile"},
    )


def google_enabled():
    return google is not None


app.jinja_env.globals["google_enabled"] = google_enabled


# ==========================================================
# 3. AI MODES CONFIGURATION
# ==========================================================
AI_MODES = {
    "developer": {
        "name": "AI Developer",
        "emoji": "\U0001f4bb",
        "tagline": "Code generation, debugging, and explanation",
        "model": "llama-3.3-70b-versatile",
        "system_prompt": (
            "You are an expert software developer and programming assistant. "
            "Write clean, well-documented code in any language requested. "
            "Debug and fix code issues with clear explanations. "
            "Explain complex programming concepts simply. "
            "Follow best practices and design patterns. "
            "You are powered by Groq and integrated into Zenith OX."
            "Always provide working code examples with proper formatting. "
            "CRITICAL FILE NAMING RULES: 1. ONLY add the File: comment to ACTUAL PROJECT FILES that the user would save. 2. Do NOT add File: to bash commands, terminal instructions, or expected output. 3. For project files, ALWAYS put the filename as a comment "
            "on the FIRST line inside the code fence. Examples:\n"
            "```python\n# File: app.py\n...code...\n```\n"
            "```html\n<!-- File: templates/index.html -->\n...code...\n```\n"
            "```css\n/* File: static/style.css */\n...code...\n```\n"
            "```javascript\n// File: static/script.js\n...code...\n```\n"
            "Use proper project paths including directories like templates/, static/, etc. "
            "CODE EXECUTION: When the user asks you to run, execute, or test code, "
            "simulate the execution by analyzing the code logically and showing the expected "
            "output in a block like:\n```output\n[Expected Output]\n```\n"
            "For bash/shell commands, simulate the terminal output similarly. "
            "Be accurate about what the code would produce."
        ),
        "temperature": 0.3,
        "max_tokens": 4000,
        "uses_web_search": False,
    },
    "story_writer": {
        "name": "AI Story Writer",
        "emoji": "\U0001f4d6",
        "tagline": "Creative writing, stories, poems, and scripts",
        "model": "llama-3.3-70b-versatile",
        "system_prompt": (
            "You are a talented creative writer. "
            "Write engaging stories with vivid descriptions and compelling characters. "
            "Craft poetry with rhythm and imagery. "
            "You are powered by Groq and integrated into Zenith OX."
            "Create scripts with authentic dialogue. "
            "Adapt your writing style to match the requested genre. "
            "Be creative, original, and evocative in your writing."
        ),
        "temperature": 0.85,
        "max_tokens": 4000,
        "uses_web_search": False,
    },
    "solve_it": {
        "name": "AI Solve It",
        "emoji": "\U0001f9ee",
        "tagline": "Math problems and step-by-step solutions",
        "model": "llama-3.3-70b-versatile",
        "system_prompt": (
            "You are an expert mathematician and problem solver. "
            "Solve math problems step-by-step, showing all work clearly. "
            "Explain mathematical concepts with examples. "
            "You are powered by Groq and integrated into Zenith OX."
            "Handle algebra, calculus, statistics, geometry, and more. "
            "Break complex problems into manageable numbered steps. "
            "Always verify your answers by checking the work."
        ),
        "temperature": 0.2,
        "max_tokens": 4000,
        "uses_web_search": False,
    },
    "researcher": {
        "name": "AI Researcher",
        "emoji": "\U0001f50d",
        "tagline": "Web search + memory research assistant",
        "model": "llama-3.3-70b-versatile",
        "system_prompt": (
            "You are Zenith OX, a secure, intelligent research assistant. "
            "You answer clearly, accurately, and concisely. "
            "You are powered by Groq and integrated into Zenith OX."
            "Use the provided past memory and web context when relevant, "
            "but never fabricate facts. If unsure, say so."
        ),
        "temperature": 0.6,
        "max_tokens": 2000,
        "uses_web_search": True,
    },
    "email_writer": {
        "name": "AI Email Writer",
        "emoji": "\u2709\ufe0f",
        "tagline": "Generate professional emails ready to copy",
        "model": "llama-3.3-70b-versatile",
        "system_prompt": (
            "You are an expert email writer. "
            "Write clear, professional, and well-structured emails. "
            "Adapt tone to context: formal, casual, follow-up, complaint, request, etc. "
            "You are powered by Groq and integrated into Zenith OX."
            "Include appropriate Subject line, greeting, body, and sign-off. "
            "Keep emails concise yet complete. "
            "Format the output as a ready-to-copy email with Subject: and Body: clearly marked."
        ),
        "temperature": 0.5,
        "max_tokens": 2000,
        "uses_web_search": False,
    },
    "pptx_generator": {
        "name": "AI Slides Generator",
        "emoji": "\U0001f4ca",
        "tagline": "Generate downloadable PowerPoint presentations",
        "model": "llama-3.3-70b-versatile",
        "system_prompt": (
            "You generate PowerPoint presentation content. "
            "You are powered by Groq and integrated into Zenith OX."
            "When the user asks for a presentation, generate ONLY a valid JSON object with no extra text.\n"
            "Use this exact format:\n"
            '{"title": "Presentation Title", "slides": ['
            '{"title": "Slide Title", "bullets": ["Point 1", "Point 2", "Point 3"]}'
            "]}\n\n"
            "Rules:\n"
            "- Generate 3-10 slides based on the topic complexity\n"
            "- Each slide should have 3-5 concise bullet points\n"
            "- First slide should be an overview\n"
            "- Last slide should be a summary or conclusion\n"
            "- Output ONLY the JSON object, no markdown fences, no explanations"
        ),
        "temperature": 0.4,
        "max_tokens": 3000,
        "uses_web_search": False,
        "special_handler": "pptx",
    },
}



# Language to file extension mapping
LANG_EXTENSIONS = {
    "python": ".py", "py": ".py",
    "javascript": ".js", "js": ".js",
    "typescript": ".ts", "ts": ".ts",
    "html": ".html", "css": ".css",
    "java": ".java", "c": ".c",
    "cpp": ".cpp", "c++": ".cpp", "csharp": ".cs",
    "go": ".go", "rust": ".rs", "ruby": ".rb",
    "php": ".php", "sql": ".sql", "swift": ".swift",
    "kotlin": ".kt", "dart": ".dart", "r": ".r",
    "bash": ".sh", "sh": ".sh", "shell": ".sh",
    "json": ".json", "yaml": ".yml", "yml": ".yml",
    "xml": ".xml", "markdown": ".md", "md": ".md",
    "txt": ".txt", "text": ".txt", "": ".txt",
}

CODE_BLOCK_RE = re.compile(r"```(\w*)\n(.*?)```", re.DOTALL)

# Allowed code file extensions for reading uploads
CODE_EXTENSIONS = {
    ".py", ".js", ".ts", ".tsx", ".jsx", ".html", ".css", ".java",
    ".c", ".cpp", ".h", ".cs", ".go", ".rs", ".rb", ".php", ".sql",
    ".swift", ".kt", ".dart", ".r", ".sh", ".json", ".yaml", ".yml",
    ".xml", ".md", ".txt", ".toml", ".cfg", ".ini", ".env", ".gitignore",
    ".dockerfile", ".makefile",
}

MAX_UPLOAD_SIZE = 10 * 1024 * 1024  # 10 MB


def extract_code_blocks(text):
    """Extract all fenced code blocks from markdown text, including filenames."""
    matches = CODE_BLOCK_RE.findall(text)
    blocks = []
    file_comment_re = re.compile(
        r"^(?:#|//|<!--|/\*)?\s*File:\s*(.+?)(?:\s*-->|\s*\*/)?\s*$",
        re.IGNORECASE
    )
    for lang, content in matches:
        content = content.strip()
        filename = None
        lines = content.split("\n")
        if lines:
            m = file_comment_re.match(lines[0].strip())
            if m:
                filename = m.group(1).strip()
                content = "\n".join(lines[1:]).strip()
        blocks.append({
            "language": (lang or "txt").lower(),
            "code": content,
            "filename": filename,
        })
    return blocks


def save_code_as_zip(code_blocks):
    """Save ONLY named code blocks as files in a zip. Skips commands & output."""
    if not code_blocks:
        return None

    # Filter: only include blocks that have a filename OR are substantial code
    named_blocks = []
    for block in code_blocks:
        lang = block["language"]
        # Always skip output/terminal blocks
        if lang in ("output", "terminal", "console"):
            continue
        # If it has a filename from the AI, always include it
        if block.get("filename"):
            named_blocks.append(block)
            continue
        # Skip short bash/shell commands (instructional, not project files)
        if lang in ("bash", "sh", "shell") and block["code"].count("\n") < 3:
            continue
        # For other languages, only include if substantial (>5 lines)
        if block["code"].count("\n") >= 5:
            ext = LANG_EXTENSIONS.get(lang, ".txt")
            block["filename"] = f"main{ext}" if len(code_blocks) == 1 else None
            if block["filename"]:
                named_blocks.append(block)

    if not named_blocks:
        return None

    timestamp = int(time.time())
    zip_filename = f"zenith_code_{timestamp}.zip"
    zip_filepath = f"/tmp/{zip_filename}"
    used_names = set()

    with zipfile.ZipFile(zip_filepath, "w", zipfile.ZIP_DEFLATED) as zf:
        for i, block in enumerate(named_blocks):
            lang = block["language"]
            ext = LANG_EXTENSIONS.get(lang, ".txt")
            fname = block.get("filename") or f"file_{i + 1}{ext}"
            # Avoid duplicate filenames
            if fname in used_names:
                base, fext = os.path.splitext(fname)
                fname = f"{base}_{i + 1}{fext}"
            used_names.add(fname)
            zf.writestr(fname, block["code"])

    if not used_names:
        return None

    return {
        "filename": zip_filename,
        "url": f"/download-zip/{zip_filename}",
        "files": sorted(used_names),
    }


def read_archive(file_storage):
    """Read code files from an uploaded zip or rar archive. Returns dict of filename->content."""
    filename = file_storage.filename.lower()
    file_bytes = file_storage.read()

    if len(file_bytes) > MAX_UPLOAD_SIZE:
        return None, "File too large (max 10 MB)"

    files = {}

    if filename.endswith(".zip"):
        try:
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
                for name in zf.namelist():
                    # Skip directories and hidden files
                    if name.endswith("/") or "/__" in name or "/." in name:
                        continue
                    ext = os.path.splitext(name)[1].lower()
                    basename = os.path.basename(name)
                    # Skip non-code files and large files
                    if ext not in CODE_EXTENSIONS and basename.lower() not in {
                        "dockerfile", "makefile", ".gitignore", ".env.example"
                    }:
                        continue
                    try:
                        content = zf.read(name).decode("utf-8", errors="replace")
                        if len(content) > 50000:  # Skip files > 50KB
                            continue
                        files[name] = content
                    except Exception:
                        continue
        except zipfile.BadZipFile:
            return None, "Invalid zip file"

    elif filename.endswith(".rar"):
        try:
            import rarfile
            with rarfile.RarFile(io.BytesIO(file_bytes)) as rf:
                for name in rf.namelist():
                    if name.endswith("/") or "/__" in name or "/." in name:
                        continue
                    ext = os.path.splitext(name)[1].lower()
                    basename = os.path.basename(name)
                    if ext not in CODE_EXTENSIONS and basename.lower() not in {
                        "dockerfile", "makefile", ".gitignore", ".env.example"
                    }:
                        continue
                    try:
                        content = rf.read(name).decode("utf-8", errors="replace")
                        if len(content) > 50000:
                            continue
                        files[name] = content
                    except Exception:
                        continue
        except ImportError:
            return None, "RAR support not available on this server. Please upload a .zip file instead."
        except Exception:
            return None, "Invalid or corrupted RAR file"
    else:
        return None, "Unsupported format. Please upload .zip or .rar"

    if not files:
        return None, "No code files found in archive"

    return files, None


# ----------------------------------------------------------
# UNIVERSAL FILE READER (supports any file type)
# ----------------------------------------------------------
# Supported file types and their extensions
SUPPORTED_FILE_TYPES = {
    # Documents
    ".pdf", ".docx", ".doc", ".txt", ".rtf", ".md",
    # Spreadsheets
    ".csv", ".xlsx", ".xls", ".tsv",
    # Code files (existing)
    ".py", ".js", ".ts", ".tsx", ".jsx", ".html", ".css", ".java",
    ".c", ".cpp", ".h", ".cs", ".go", ".rs", ".rb", ".php", ".sql",
    ".swift", ".kt", ".dart", ".r", ".sh", ".json", ".yaml", ".yml",
    ".xml", ".toml", ".cfg", ".ini", ".env", ".gitignore",
    ".dockerfile", ".makefile",
    # Images (text description only)
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".svg",
    # Archives
    ".zip", ".rar",
    # Presentations
    ".pptx", ".ppt",
}


def extract_file_content(file_storage):
    """
    Extract text content from ANY uploaded file.
    Returns (content_string, error_string).
    For unsupported types, returns a helpful message.
    """
    filename = file_storage.filename or "unknown"
    ext = os.path.splitext(filename)[1].lower()
    file_bytes = file_storage.read()

    if len(file_bytes) > MAX_UPLOAD_SIZE:
        return None, "File too large (max 10 MB)"

    try:
        # === PDF FILES ===
        if ext == ".pdf":
            try:
                import pdfplumber
                text_parts = []
                with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                    for i, page in enumerate(pdf.pages[:50]):  # Max 50 pages
                        page_text = page.extract_text() or ""
                        if page_text.strip():
                            text_parts.append(f"--- Page {i+1} ---\n{page_text}")
                if text_parts:
                    content = "\n\n".join(text_parts)
                    if len(content) > 30000:
                        content = content[:30000] + "\n\n[... truncated due to length ...]"
                    return content, None
                return None, "Could not extract text from PDF (may be image-based)"
            except ImportError:
                return None, "PDF support requires: pip install pdfplumber"

        # === WORD DOCUMENTS ===
        elif ext in (".docx", ".doc"):
            try:
                from docx import Document
                doc = Document(io.BytesIO(file_bytes))
                paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                if paragraphs:
                    content = "\n\n".join(paragraphs)
                    if len(content) > 30000:
                        content = content[:30000] + "\n\n[... truncated ...]"
                    return content, None
                return None, "Word document appears to be empty"
            except ImportError:
                return None, "DOCX support requires: pip install python-docx"

        # === EXCEL / SPREADSHEETS ===
        elif ext in (".xlsx", ".xls"):
            try:
                import openpyxl
                wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True)
                sheets_text = []
                for sheet_name in wb.sheetnames[:10]:  # Max 10 sheets
                    ws = wb[sheet_name]
                    rows = []
                    for row in ws.iter_rows(max_row=200, values_only=True):  # Max 200 rows
                        row_str = " | ".join(str(cell) if cell is not None else "" for cell in row)
                        if row_str.strip(" |"):
                            rows.append(row_str)
                    if rows:
                        sheets_text.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
                wb.close()
                if sheets_text:
                    content = "\n\n".join(sheets_text)
                    if len(content) > 30000:
                        content = content[:30000] + "\n\n[... truncated ...]"
                    return content, None
                return None, "Excel file appears to be empty"
            except ImportError:
                return None, "Excel support requires: pip install openpyxl"

        # === CSV / TSV ===
        elif ext in (".csv", ".tsv"):
            try:
                text = file_bytes.decode("utf-8", errors="replace")
                delimiter = "\t" if ext == ".tsv" else ","
                reader = csv.reader(text.splitlines(), delimiter=delimiter)
                rows = []
                for i, row in enumerate(reader):
                    if i >= 500:  # Max 500 rows
                        rows.append("[... more rows truncated ...]")
                        break
                    rows.append(" | ".join(row))
                if rows:
                    content = "\n".join(rows)
                    if len(content) > 30000:
                        content = content[:30000] + "\n\n[... truncated ...]"
                    return content, None
                return None, "CSV file appears to be empty"
            except Exception:
                return None, "Could not parse CSV file"

        # === POWERPOINT ===
        elif ext in (".pptx", ".ppt"):
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(file_bytes))
                slides_text = []
                for i, slide in enumerate(prs.slides):
                    texts = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            texts.append(shape.text)
                    if texts:
                        slides_text.append(f"--- Slide {i+1} ---\n" + "\n".join(texts))
                if slides_text:
                    return "\n\n".join(slides_text), None
                return None, "PowerPoint file appears to be empty"
            except ImportError:
                return None, "PPTX support requires: pip install python-pptx"

        # === IMAGES (describe what was uploaded) ===
        elif ext in (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"):
            try:
                from PIL import Image
                img = Image.open(io.BytesIO(file_bytes))
                w, h = img.size
                mode = img.mode
                fmt = img.format or ext.upper()
                desc = (
                    f"[Uploaded Image: {filename}]\n"
                    f"Format: {fmt}, Size: {w}x{h} pixels, Color mode: {mode}\n"
                    f"Note: This is an image file. I can discuss it based on context, "
                    f"but I cannot visually analyze the image content directly."
                )
                return desc, None
            except ImportError:
                return f"[Uploaded Image: {filename}] (Install Pillow for metadata)", None
            except Exception:
                return f"[Uploaded Image: {filename}]", None

        # === SVG ===
        elif ext == ".svg":
            text = file_bytes.decode("utf-8", errors="replace")
            if len(text) > 30000:
                text = text[:30000] + "\n[... truncated ...]"
            return f"[SVG Image: {filename}]\n{text}", None

        # === ARCHIVES (use existing handler) ===
        elif ext in (".zip", ".rar"):
            file_storage.seek(0)
            files, err = read_archive(file_storage)
            if err:
                return None, err
            return format_files_for_prompt(files), None

        # === PLAIN TEXT / CODE FILES ===
        elif ext in CODE_EXTENSIONS or ext in (".txt", ".md", ".rtf"):
            text = file_bytes.decode("utf-8", errors="replace")
            if len(text) > 50000:
                text = text[:50000] + "\n\n[... truncated ...]"
            return text, None

        # === UNKNOWN FILE TYPE ===
        else:
            # Try to read as text anyway
            try:
                text = file_bytes.decode("utf-8", errors="strict")
                if len(text) > 30000:
                    text = text[:30000] + "\n\n[... truncated ...]"
                return f"[File: {filename}]\n{text}", None
            except UnicodeDecodeError:
                return None, f"Unsupported file type: {ext}. I can analyze text documents, PDFs, Word files, spreadsheets, code files, and more."
    except Exception as e:
        return None, f"Error reading file: {str(e)}"


def format_files_for_prompt(files_dict):
    """Format extracted files into a string for the AI prompt."""
    parts = []
    for filepath, content in sorted(files_dict.items()):
        parts.append(f"--- {filepath} ---\n{content}")
    combined = "\n\n".join(parts)
    # Truncate if too long (keep under ~30K chars for the prompt)
    if len(combined) > 30000:
        combined = combined[:30000] + "\n\n[... truncated due to length ...]"
    return combined


# ==========================================================
# 4. STORAGE HELPERS
# ==========================================================
def load_users():
    for path in [WRITABLE_USERS, USERS_FILE]:
        if os.path.exists(path):
            try:
                with open(path, "r", encoding="utf-8") as f:
                    return json.load(f)
            except (json.JSONDecodeError, OSError):
                continue
    return {}


def save_users(users):
    with open(WRITABLE_USERS, "w", encoding="utf-8") as f:
        json.dump(users, f, indent=2, ensure_ascii=False)


def load_memory():
    for path in [WRITABLE_MEMORY, MEMORY_FILE]:
        if os.path.exists(path):
            try:
                with open(path, "r", encoding="utf-8") as f:
                    return json.load(f)
            except (json.JSONDecodeError, OSError):
                continue
    return {}


def save_memory(memory):
    with open(WRITABLE_MEMORY, "w", encoding="utf-8") as f:
        json.dump(memory, f, indent=2, ensure_ascii=False)


def get_user_memory(memory_key):
    return load_memory().get(memory_key, [])


def update_user_memory(memory_key, role, content):
    mem = load_memory()
    mem.setdefault(memory_key, [])
    mem[memory_key].append({"role": role, "content": content})
    if len(mem[memory_key]) > MEMORY_LIMIT * 2:
        mem[memory_key] = mem[memory_key][-MEMORY_LIMIT * 2:]
    save_memory(mem)


# ==========================================================
# 4b. AUTH HELPERS
# ==========================================================
EMAIL_RE = re.compile(r"^[^@\s]+@[^@\s]+\.[^@\s]+$")


def valid_email(email: str) -> bool:
    return bool(EMAIL_RE.match(email or ""))


def find_user_by_email(email: str):
    email = (email or "").strip().lower()
    users = load_users()
    for key, u in users.items():
        if isinstance(u, dict) and u.get("email", "").lower() == email:
            return key, u
    return None, None


def find_user_by_google_id(gid: str):
    users = load_users()
    for key, u in users.items():
        if isinstance(u, dict) and u.get("google_id") == gid:
            return key, u
    return None, None


def display_name_from_email(email: str) -> str:
    return (email or "").split("@")[0] or "friend"


# ==========================================================
# 4c. TIME-BASED GREETING
# ==========================================================
def time_based_greeting(name: str) -> str:
    hour = datetime.now().hour
    if 5 <= hour < 12:
        part = "Good morning"
    elif 12 <= hour < 17:
        part = "Good afternoon"
    elif 17 <= hour < 22:
        part = "Good evening"
    else:
        part = "Good night"
    return f"{part}, {name}"


# ==========================================================
# 5. VECTOR MEMORY (TF-IDF + cosine similarity)
# ==========================================================
def retrieve_relevant_memory(memory_key, query, top_k=TOP_K_MEMORY):
    history = get_user_memory(memory_key)
    if not history:
        return ""

    pairs = []
    for i, msg in enumerate(history):
        if msg["role"] == "user":
            reply = ""
            if i + 1 < len(history) and history[i + 1]["role"] == "assistant":
                reply = history[i + 1]["content"]
            pairs.append((msg["content"], reply))

    if not pairs:
        return ""

    user_texts = [p[0] for p in pairs]

    try:
        vectorizer = TfidfVectorizer(stop_words="english")
        matrix = vectorizer.fit_transform(user_texts + [query])
        query_vec = matrix[-1]
        past_vecs = matrix[:-1]
        sims = cosine_similarity(query_vec, past_vecs).flatten()
    except ValueError:
        return ""

    ranked = np.argsort(sims)[::-1]
    selected = [idx for idx in ranked if sims[idx] > 0][:top_k]

    if not selected:
        return ""

    lines = []
    for idx in selected:
        u, a = pairs[idx]
        lines.append(f"- User previously asked: {u}\n  You answered: {a}")
    return "\n".join(lines)


# ==========================================================
# 6. WEB SEARCH (Tavily)
# ==========================================================
def tavily_search(query, max_results=3):
    if not TAVILY_API_KEY:
        return ""
    try:
        r = requests.post(
            "https://api.tavily.com/search",
            json={
                "api_key": TAVILY_API_KEY,
                "query": query,
                "search_depth": "basic",
                "max_results": max_results,
            },
            timeout=15,
        )
        r.raise_for_status()
        data = r.json()
        chunks = [res.get("content", "") for res in data.get("results", [])]
        return "\n\n".join(filter(None, chunks))[:4000]
    except Exception as e:
        print("[Tavily error]", e)
        return ""


# ==========================================================
# 7. GROQ CHAT (mode-aware)
# Vision model for image analysis
VISION_MODEL = "llama-4-scout-17b-16e-instruct"

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp"}


def ask_groq_vision(user_input, image_bytes, image_ext, mode, recent_history=None):
    """Send an image + text prompt to Groq's vision model."""
    # Determine MIME type
    mime_map = {
        ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
        ".gif": "image/gif", ".bmp": "image/bmp", ".webp": "image/webp",
    }
    mime = mime_map.get(image_ext.lower(), "image/jpeg")

    # Encode image as base64
    b64 = base64.b64encode(image_bytes).decode("utf-8")
    data_url = f"data:{mime};base64,{b64}"

    # Build messages
    messages = [{"role": "system", "content": mode["system_prompt"]}]
    if recent_history:
        for msg in recent_history[-6:]:
            messages.append({"role": msg["role"], "content": msg["content"]})

    # Multimodal user message
    messages.append({
        "role": "user",
        "content": [
            {"type": "text", "text": user_input or "Please analyze this image."},
            {"type": "image_url", "image_url": {"url": data_url}},
        ],
    })

    try:
        resp = groq_client.chat.completions.create(
            model=VISION_MODEL,
            messages=messages,
            temperature=mode["temperature"],
            max_tokens=mode["max_tokens"],
        )
        return resp.choices[0].message.content.strip()
    except Exception as e:
        return f"[Zenith OX error] {e}"
# ==========================================================
def ask_groq(user_input, vector_memory, web_context, mode, recent_history=None):
    # Build messages array with conversation history for context
    messages = [{"role": "system", "content": mode["system_prompt"]}]

    # Include recent conversation history so the model remembers context
    if recent_history:
        for msg in recent_history[-10:]:  # Last 5 exchanges (10 messages)
            messages.append({"role": msg["role"], "content": msg["content"]})

    # Build current user prompt with memory/web context
    prompt = f"User Question: {user_input}\n\n"
    if vector_memory:
        prompt += f"Relevant Past Memory:\n{vector_memory}\n\n"
    if web_context:
        prompt += f"Web Context:\n{web_context}\n\n"
    if not mode.get("special_handler"):
        prompt += "Instruction:\nProvide a clear, accurate, and helpful answer."

    messages.append({"role": "user", "content": prompt})

    try:
        resp = groq_client.chat.completions.create(
            model=mode["model"],
            messages=messages,
            temperature=mode["temperature"],
            max_tokens=mode["max_tokens"],
        )
        return resp.choices[0].message.content.strip()
    except Exception as e:
        return f"[Zenith OX error] {e}"


# ==========================================================
# 8. PPTX GENERATOR
# ==========================================================
def generate_pptx(ai_response, user_id):
    """Parse AI JSON response and create a .pptx file."""
    try:
        text = ai_response.strip()
        if text.startswith("```"):
            text = text.split("\n", 1)[1]
            text = text.rsplit("```", 1)[0].strip()
        data = json.loads(text)
    except (json.JSONDecodeError, IndexError):
        return None

    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    title = data.get("title", "Presentation")
    slides = data.get("slides", [])

    if not slides:
        return None

    # Title slide
    title_layout = prs.slide_layouts[0]
    first = prs.slides.add_slide(title_layout)
    first.shapes.title.text = title
    if len(first.placeholders) > 1:
        first.placeholders[1].text = "Generated by Zenith OX"

    # Content slides
    content_layout = prs.slide_layouts[1]
    slide_titles = []
    for s in slides:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = s.get("title", "Untitled")
        slide_titles.append(s.get("title", "Untitled"))

        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        for j, bullet in enumerate(s.get("bullets", [])):
            if j == 0:
                tf.text = bullet
            else:
                p = tf.add_paragraph()
                p.text = bullet

    filename = f"zenith_ox_{int(time.time())}.pptx"
    filepath = f"/tmp/{filename}"
    prs.save(filepath)

    return {
        "filename": filename,
        "path": filepath,
        "url": f"/download/{filename}",
        "slides": slide_titles,
    }


# ==========================================================
# 9. ROUTES
# ==========================================================

# ---------- INDEX ----------
@app.route("/")
def index():
    if "user_id" not in session:
        return redirect(url_for("login_page"))
    if "ai_mode" not in session:
        return redirect(url_for("menu"))

    mode_key = session["ai_mode"]
    mode = AI_MODES.get(mode_key, AI_MODES["researcher"])
    username = session.get("display_name") or session["user_id"]
    greeting = time_based_greeting(username)
    return render_template(
        "index.html",
        username=username,
        greeting=greeting,
        mode=mode,
        mode_key=mode_key,
    )


# ---------- MENU ----------
@app.route("/menu")
def menu():
    if "user_id" not in session:
        return redirect(url_for("login_page"))
    username = session.get("display_name") or session["user_id"]
    greeting = time_based_greeting(username)
    return render_template(
        "menu.html",
        username=username,
        greeting=greeting,
        modes=AI_MODES,
    )


# ---------- SELECT MODE ----------
@app.route("/select-mode/<mode_key>")
def select_mode(mode_key):
    if "user_id" not in session:
        return redirect(url_for("login_page"))
    if mode_key not in AI_MODES:
        flash("Invalid AI mode.", "error")
        return redirect(url_for("menu"))
    session["ai_mode"] = mode_key
    return redirect(url_for("index"))


# ---------- REGISTER ----------
@app.route("/register", methods=["GET", "POST"])
def register():
    if request.method == "GET":
        return render_template("register.html")

    email    = (request.form.get("email") or "").strip().lower()
    password = request.form.get("password") or ""
    confirm  = request.form.get("confirm") or ""
    name     = (request.form.get("name") or "").strip()

    if not valid_email(email):
        flash("Please enter a valid email address.", "error")
        return redirect(url_for("register"))

    if len(password) < 6:
        flash("Password must be at least 6 characters.", "error")
        return redirect(url_for("register"))

    if password != confirm:
        flash("Passwords do not match.", "error")
        return redirect(url_for("register"))

    users = load_users()
    existing_key, _ = find_user_by_email(email)
    if existing_key:
        flash("An account with that email already exists. Please log in.", "error")
        return redirect(url_for("login_page"))

    user_key = email
    users[user_key] = {
        "email": email,
        "name": name or display_name_from_email(email),
        "password_hash": generate_password_hash(password),
        "google_id": None,
        "created_at": datetime.utcnow().isoformat() + "Z",
    }
    save_users(users)

    session["user_id"]      = user_key
    session["display_name"] = users[user_key]["name"]
    flash("Account created. Welcome!", "success")
    return redirect(url_for("menu"))


# ---------- LOGIN ----------
@app.route("/login", methods=["GET", "POST"])
def login_page():
    if request.method == "GET":
        return render_template("login.html")

    email    = (request.form.get("email") or "").strip().lower()
    password = request.form.get("password") or ""

    if not valid_email(email) or not password:
        flash("Please enter your email and password.", "error")
        return redirect(url_for("login_page"))

    key, user = find_user_by_email(email)
    if not user or not user.get("password_hash"):
        flash("No account found with that email, or it was created via Google. "
              "Try \'Continue with Google\' instead.", "error")
        return redirect(url_for("login_page"))

    if not check_password_hash(user["password_hash"], password):
        flash("Incorrect password.", "error")
        return redirect(url_for("login_page"))

    session["user_id"]      = key
    session["display_name"] = user.get("name") or display_name_from_email(email)
    return redirect(url_for("menu"))


# ---------- GOOGLE OAUTH ----------
@app.route("/login/google")
def login_google():
    if not google:
        flash("Google login is not configured on this server.", "error")
        return redirect(url_for("login_page"))
    redirect_uri = url_for("auth_google_callback", _external=True)
    return google.authorize_redirect(redirect_uri)


@app.route("/auth/google/callback")
def auth_google_callback():
    if not google:
        flash("Google login is not configured on this server.", "error")
        return redirect(url_for("login_page"))

    try:
        token = google.authorize_access_token()
    except Exception as e:
        flash(f"Google sign-in failed: {e}", "error")
        return redirect(url_for("login_page"))

    userinfo = token.get("userinfo") or {}
    if not userinfo:
        try:
            userinfo = google.parse_id_token(token) or {}
        except Exception:
            userinfo = {}

    google_id = userinfo.get("sub")
    email     = (userinfo.get("email") or "").strip().lower()
    name      = userinfo.get("name") or display_name_from_email(email)

    if not google_id or not email:
        flash("Google did not return the required profile info.", "error")
        return redirect(url_for("login_page"))

    users = load_users()
    key, user = find_user_by_google_id(google_id)

    if not user:
        key, user = find_user_by_email(email)
        if user:
            user["google_id"] = google_id
            user.setdefault("name", name)
            users[key] = user
            save_users(users)

    if not user:
        key = email
        users[key] = {
            "email": email,
            "name": name,
            "password_hash": None,
            "google_id": google_id,
            "created_at": datetime.utcnow().isoformat() + "Z",
        }
        save_users(users)
        user = users[key]

    session["user_id"]      = key
    session["display_name"] = user.get("name") or display_name_from_email(email)
    return redirect(url_for("menu"))


# ---------- LOGOUT ----------
@app.route("/logout", methods=["POST", "GET"])
def logout():
    session.pop("user_id", None)
    session.pop("display_name", None)
    session.pop("ai_mode", None)
    return redirect(url_for("login_page"))


# ---------- CHAT ----------
# Keywords that indicate user wants file generation/download
FILE_GEN_KEYWORDS = [
    "create a file", "generate a file", "make a file", "write a file",
    "save as", "download", "export", "create file", "generate file",
    "make file", "write file", "give me the file", "give me a file",
    "create the code", "generate the code", "zip", "package",
    "save the code", "downloadable", "create a project",
    "generate a project", "make a project", "build a project",
]


def user_wants_file(message):
    """Check if the user's message indicates they want a downloadable file."""
    msg_lower = message.lower()
    return any(kw in msg_lower for kw in FILE_GEN_KEYWORDS)


@app.route("/chat", methods=["POST"])
def chat():
    if "user_id" not in session:
        return jsonify({"ok": False, "error": "Not authenticated."}), 401

    user_id  = session["user_id"]
    mode_key = session.get("ai_mode", "researcher")
    mode     = AI_MODES.get(mode_key, AI_MODES["researcher"])

    # Support both JSON and multipart/form-data (for file uploads)
    file_context = ""
    uploaded_file = None

    if request.content_type and "multipart/form-data" in request.content_type:
        message = (request.form.get("message") or "").strip()
        uploaded_file = request.files.get("file")
    else:
        data    = request.get_json(silent=True) or {}
        message = (data.get("message") or "").strip()

    if not message and not uploaded_file:
        return jsonify({"ok": False, "error": "Empty message."}), 400

    # ── Handle file upload (any file type) ──
    is_image_upload = False
    image_bytes = None
    image_ext = None

    if uploaded_file and uploaded_file.filename:
        ext = os.path.splitext(uploaded_file.filename)[1].lower()

        # Check if it's an image → use vision model
        if ext in IMAGE_EXTENSIONS:
            is_image_upload = True
            image_bytes = uploaded_file.read()
            image_ext = ext
            if len(image_bytes) > MAX_UPLOAD_SIZE:
                return jsonify({"ok": False, "error": "File too large (max 10 MB)"}), 400
        else:
            # Non-image file: extract text content
            content, error = extract_file_content(uploaded_file)
            if error:
                return jsonify({"ok": False, "error": error}), 400
            if content:
                file_context = (
                    f"\n\n--- Uploaded File: {uploaded_file.filename} ---\n"
                    f"{content}\n--- End of File ---\n"
                )

    memory_key = f"{user_id}:{mode_key}"
    recent_history = get_user_memory(memory_key)

    # ── Route to vision model for images, text model for everything else ──
    if is_image_upload:
        answer = ask_groq_vision(
            message or "Please analyze this image.",
            image_bytes, image_ext, mode,
            recent_history=recent_history,
        )
        full_message = message or f"[Uploaded image: {uploaded_file.filename}]"
    else:
        # Combine message with file context
        full_message = message
        if file_context:
            full_message = f"{message}\n{file_context}" if message else f"Please analyze this file:{file_context}"

        if not full_message.strip():
            return jsonify({"ok": False, "error": "Empty message."}), 400

        vector_mem = retrieve_relevant_memory(memory_key, message or "file analysis")
        web_ctx    = tavily_search(message) if mode.get("uses_web_search") and message else ""
        answer     = ask_groq(full_message, vector_mem, web_ctx, mode, recent_history=recent_history)

    # Handle PPTX special mode — always generates file (that's its purpose)
    if mode.get("special_handler") == "pptx":
        try:
            result = generate_pptx(answer, user_id)
            if result:
                summary = f"Your presentation is ready!\n\nSlides:\n"
                for i, t in enumerate(result["slides"], 1):
                    summary += f"  {i}. {t}\n"
                summary += "\nClick the download button below to save your file."

                update_user_memory(memory_key, "user", message)
                update_user_memory(memory_key, "assistant", summary)
                return jsonify({
                    "ok": True,
                    "response": summary,
                    "download_url": result["url"],
                    "download_name": result["filename"],
                })
        except Exception as e:
            print(f"[PPTX error] {e}")

    # For developer mode: only generate zip when user explicitly asks for files
    if mode_key == "developer" and user_wants_file(message):
        code_blocks = extract_code_blocks(answer)
        if code_blocks:
            zip_info = save_code_as_zip(code_blocks)
            if zip_info:
                update_user_memory(memory_key, "user", message)
                update_user_memory(memory_key, "assistant", answer)
                return jsonify({
                    "ok": True,
                    "response": answer,
                    "download_url": zip_info["url"],
                    "download_name": zip_info["filename"],
                })

    # Store in memory (store original message, not file content, to keep memory clean)
    update_user_memory(memory_key, "user", message or f"[Uploaded: {uploaded_file.filename}]")
    update_user_memory(memory_key, "assistant", answer)
    return jsonify({"ok": True, "response": answer})


# ---------- CLEAR MEMORY ----------
@app.route("/clear", methods=["POST"])
def clear():
    if "user_id" not in session:
        return jsonify({"ok": False, "error": "Not authenticated."}), 401

    user_id    = session["user_id"]
    mode_key   = session.get("ai_mode", "researcher")
    memory_key = f"{user_id}:{mode_key}"

    mem = load_memory()
    if memory_key in mem:
        mem[memory_key] = []
        save_memory(mem)
    return jsonify({"ok": True, "message": "Memory cleared for this mode."})


# ---------- DOWNLOAD (PPTX) ----------
@app.route("/download/<filename>")
def download_file(filename):
    if "user_id" not in session:
        return redirect(url_for("login_page"))
    if not filename.endswith(".pptx"):
        return "Invalid file type", 400
    filepath = os.path.join("/tmp", filename)
    if not os.path.exists(filepath):
        return "File not found or expired", 404
    return send_from_directory("/tmp", filename, as_attachment=True)



# ---------- UPLOAD CODE (ZIP/RAR) ----------
@app.route("/upload-code", methods=["POST"])
def upload_code():
    """Handle zip/rar upload: read files, send to AI for modification, return new zip."""
    if "user_id" not in session:
        return jsonify({"ok": False, "error": "Not authenticated."}), 401

    if "file" not in request.files:
        return jsonify({"ok": False, "error": "No file uploaded."}), 400

    uploaded = request.files["file"]
    if not uploaded.filename:
        return jsonify({"ok": False, "error": "Empty filename."}), 400

    instruction = request.form.get("message", "").strip()
    if not instruction:
        instruction = "Analyze this code and suggest improvements."

    # Read the archive
    files_dict, error = read_archive(uploaded)
    if error:
        return jsonify({"ok": False, "error": error}), 400

    # Format files for the AI prompt
    files_text = format_files_for_prompt(files_dict)
    file_list = ", ".join(files_dict.keys())

    user_id  = session["user_id"]
    mode_key = "developer"
    mode     = AI_MODES["developer"]
    memory_key = f"{user_id}:{mode_key}"

    # Build a rich prompt with the uploaded code
    upload_prompt = (
        f"The user uploaded a code project with these files: {file_list}\n\n"
        f"Here are the file contents:\n\n{files_text}\n\n"
        f"User instruction: {instruction}\n\n"
        "IMPORTANT: When returning modified code, wrap each file in markdown code fences "
        "with the language name. If modifying multiple files, include ALL of them."
    )

    recent_history = get_user_memory(memory_key)
    vector_mem = retrieve_relevant_memory(memory_key, instruction)
    answer = ask_groq(upload_prompt, vector_mem, "", mode, recent_history=recent_history)

    # Extract code blocks from the response and create a zip
    code_blocks = extract_code_blocks(answer)
    response_data = {"ok": True, "response": answer}

    if code_blocks:
        zip_info = save_code_as_zip(code_blocks)
        if zip_info:
            response_data["download_url"] = zip_info["url"]
            response_data["download_name"] = zip_info["filename"]

    update_user_memory(memory_key, "user", f"[Uploaded: {file_list}] {instruction}")
    update_user_memory(memory_key, "assistant", answer)

    return jsonify(response_data)


# ---------- DOWNLOAD ZIP ----------
@app.route("/download-zip/<filename>")
def download_zip(filename):
    if "user_id" not in session:
        return redirect(url_for("login_page"))
    if not filename.endswith(".zip"):
        return "Invalid file type", 400
    filepath = os.path.join("/tmp", filename)
    if not os.path.exists(filepath):
        return "File not found or expired", 404
    return send_from_directory("/tmp", filename, as_attachment=True)



# ---------- CHAT HISTORY (for loading saved chats) ----------
@app.route("/history")
def history():
    if "user_id" not in session:
        return jsonify({"ok": False, "error": "Not authenticated."}), 401
    user_id  = session["user_id"]
    mode_key = session.get("ai_mode", "researcher")
    memory_key = f"{user_id}:{mode_key}"
    messages = get_user_memory(memory_key)
    # Return last 30 messages (15 exchanges)
    return jsonify({"ok": True, "messages": messages[-30:]})


# ==========================================================
# 10. ENTRY POINT
# ==========================================================
if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000, debug=True)

application = app
